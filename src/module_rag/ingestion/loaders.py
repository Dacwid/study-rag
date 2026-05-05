"""Document loaders for PDF, PPTX, DOCX, and Markdown/text formats.

Each loader returns a list of LangChain Documents with rich metadata so that
downstream retrieval can filter by module, week, or content type without
re-parsing the original file.
"""

from __future__ import annotations

import re
from pathlib import Path

from langchain_core.documents import Document

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".md", ".txt"}


def _parse_module(path: Path) -> str:
    """Infer module name from the immediate parent folder, or filename stem."""
    parent = path.parent.name
    if parent and parent.lower() not in ("raw", ".", ""):
        return parent
    match = re.match(r"([a-zA-Z]{2,}\d+)", path.stem)
    return match.group(1) if match else "unknown"


def _parse_week(path: Path) -> int | None:
    """Extract week number from filename, e.g. 'week3_slides.pdf' → 3."""
    match = re.search(r"week[_\-]?(\d+)", path.stem, re.IGNORECASE)
    return int(match.group(1)) if match else None


def load_pdf(path: Path) -> list[Document]:
    """Load a PDF — one Document per page.

    Falls back to pytesseract OCR when a page yields fewer than 50 characters
    of extractable text (common for scanned/image-only slides).

    Args:
        path: Path to the PDF file.

    Returns:
        List of Documents, one per non-empty page.
    """
    import fitz

    module = _parse_module(path)
    week = _parse_week(path)
    docs: list[Document] = []

    with fitz.open(str(path)) as pdf:
        for page_num, page in enumerate(pdf, start=1):
            text = page.get_text().strip()

            if len(text) < 50:
                text = _ocr_page(page)

            if not text:
                continue

            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source_file": str(path),
                        "module": module,
                        "week": week,
                        "page_or_slide": page_num,
                        "content_type": "pdf_page",
                    },
                )
            )

    return docs


def _ocr_page(page: object) -> str:
    """Render a pymupdf page and run Tesseract OCR on it."""
    try:
        import io

        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=200)  # type: ignore[attr-defined]
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ""


def load_pptx(path: Path) -> list[Document]:
    """Load a PPTX — separate Documents for title, body, and notes per slide.

    Splitting by field preserves structure so title text is never mixed with
    body bullets in the same embedding, improving retrieval precision.

    Args:
        path: Path to the PPTX file.

    Returns:
        List of Documents (up to 3 per slide: title, body, notes).
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    module = _parse_module(path)
    week = _parse_week(path)
    docs: list[Document] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        base_meta = {
            "source_file": str(path),
            "module": module,
            "week": week,
            "page_or_slide": slide_num,
        }

        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
            if title_text:
                docs.append(
                    Document(
                        page_content=title_text,
                        metadata={**base_meta, "content_type": "slide_title"},
                    )
                )

        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)
        if body_parts:
            docs.append(
                Document(
                    page_content="\n".join(body_parts),
                    metadata={**base_meta, "content_type": "slide_body"},
                )
            )

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                docs.append(
                    Document(
                        page_content=notes_text,
                        metadata={**base_meta, "content_type": "slide_notes"},
                    )
                )

    return docs


def load_docx(path: Path) -> list[Document]:
    """Load a DOCX — one Document per paragraph, preserving heading metadata.

    Args:
        path: Path to the DOCX file.

    Returns:
        List of Documents, one per non-empty paragraph.
    """
    from docx import Document as DocxDocument

    docx = DocxDocument(str(path))
    module = _parse_module(path)
    week = _parse_week(path)
    docs: list[Document] = []

    for para_num, para in enumerate(docx.paragraphs, start=1):
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""
        heading_match = re.match(r"Heading (\d+)", style_name)
        meta: dict = {
            "source_file": str(path),
            "module": module,
            "week": week,
            "page_or_slide": para_num,
            "content_type": "docx_paragraph",
        }
        if heading_match:
            meta["heading_level"] = int(heading_match.group(1))

        docs.append(Document(page_content=text, metadata=meta))

    return docs


def load_text(path: Path) -> list[Document]:
    """Load a Markdown or plain-text file.

    Splits on H2 headings (``## ...``) when present; otherwise treats the
    whole file as a single Document.

    Args:
        path: Path to the .md or .txt file.

    Returns:
        List of Documents, one per H2 section (or one for the whole file).
    """
    text = path.read_text(encoding="utf-8")
    module = _parse_module(path)
    week = _parse_week(path)

    sections = re.split(r"(?m)^(?=## )", text)
    sections = [s.strip() for s in sections if s.strip()]

    if not sections:
        return []

    if len(sections) == 1:
        return [
            Document(
                page_content=sections[0],
                metadata={
                    "source_file": str(path),
                    "module": module,
                    "week": week,
                    "page_or_slide": 1,
                    "content_type": "text",
                },
            )
        ]

    return [
        Document(
            page_content=section,
            metadata={
                "source_file": str(path),
                "module": module,
                "week": week,
                "page_or_slide": idx,
                "content_type": "text",
            },
        )
        for idx, section in enumerate(sections, start=1)
    ]


_LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".md": load_text,
    ".txt": load_text,
}


def load_file(path: Path) -> list[Document]:
    """Dispatch to the correct loader based on file extension.

    Args:
        path: Path to a supported document.

    Returns:
        List of Documents with rich metadata.

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = path.suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return loader(path)
